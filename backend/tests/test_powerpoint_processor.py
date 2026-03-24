"""
Tests for PowerPoint Document Processor

Tests the PowerPointProcessor class for text extraction, format preservation,
and file validation.
"""

import pytest
import tempfile
from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from src.services.document_processor import DocumentType
from src.services.powerpoint_processor import PowerPointProcessor


class TestPowerPointProcessor:
    """Tests for PowerPointProcessor."""
    
    def test_supported_extensions(self):
        """Test supported extensions."""
        processor = PowerPointProcessor()
        assert processor.supported_extensions == ['.pptx']
    
    def test_document_type(self):
        """Test document type."""
        processor = PowerPointProcessor()
        assert processor.document_type == DocumentType.POWERPOINT
    
    def test_generate_output_filename(self):
        """Test output filename generation."""
        processor = PowerPointProcessor()
        
        filename = processor.generate_output_filename(Path("presentation.pptx"))
        assert filename == "presentation_vi.pptx"
        
        filename = processor.generate_output_filename(Path("slides.pptx"), "en")
        assert filename == "slides_en.pptx"


class TestPowerPointProcessorExtraction:
    """Tests for PowerPointProcessor text extraction."""
    
    @pytest.mark.asyncio
    async def test_extract_text_from_shapes(self):
        """Test extracting text from slide shapes."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add a text box
            left = top = Pt(100)
            width = height = Pt(200)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = "Hello World"
            
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(file_path)
            
            assert len(segments) >= 1
            texts = [s.text for s in segments]
            assert "Hello World" in texts
            
            # Check metadata
            shape_segments = [s for s in segments if s.metadata.get("type") == "shape"]
            assert len(shape_segments) >= 1
        finally:
            file_path.unlink()


    @pytest.mark.asyncio
    async def test_extract_text_from_tables(self):
        """Test extracting text from tables on slides."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add a table
            left = top = Pt(100)
            width = Pt(400)
            height = Pt(200)
            table = slide.shapes.add_table(2, 2, left, top, width, height).table
            table.cell(0, 0).text = "Cell A1"
            table.cell(0, 1).text = "Cell B1"
            table.cell(1, 0).text = "Cell A2"
            table.cell(1, 1).text = "Cell B2"
            
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(file_path)
            
            # Should have 4 table cell segments
            table_segments = [s for s in segments if s.metadata.get("type") == "table_cell"]
            assert len(table_segments) == 4
            
            texts = [s.text for s in table_segments]
            assert "Cell A1" in texts
            assert "Cell B1" in texts
            assert "Cell A2" in texts
            assert "Cell B2" in texts
        finally:
            file_path.unlink()
    
    @pytest.mark.asyncio
    async def test_extract_text_from_notes(self):
        """Test extracting text from speaker notes."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Speaker notes content"
            
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(file_path)
            
            # Should have notes segment
            notes_segments = [s for s in segments if s.metadata.get("type") == "notes"]
            assert len(notes_segments) >= 1
            
            texts = [s.text for s in notes_segments]
            assert "Speaker notes content" in texts
        finally:
            file_path.unlink()
    
    @pytest.mark.asyncio
    async def test_extract_text_preserves_formatting_metadata(self):
        """Test that extraction preserves formatting metadata."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add a text box with formatted text
            txBox = slide.shapes.add_textbox(Pt(100), Pt(100), Pt(200), Pt(100))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Bold Text"
            run.font.bold = True
            run.font.size = Pt(14)
            
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(file_path)
            
            assert len(segments) >= 1
            segment = [s for s in segments if "Bold Text" in s.text][0]
            
            # Check runs metadata
            runs = segment.metadata.get("runs", [])
            assert len(runs) >= 1
            assert runs[0]["bold"] is True
            assert runs[0]["font_size"] == 14.0
        finally:
            file_path.unlink()


class TestPowerPointProcessorValidation:
    """Tests for PowerPointProcessor file validation."""
    
    @pytest.mark.asyncio
    async def test_validate_file_success(self):
        """Test validating a valid PowerPoint file."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Pt(100), Pt(100), Pt(200), Pt(100))
            txBox.text_frame.text = "Test content"
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            is_valid, error = await processor.validate_file(file_path)
            assert is_valid is True
            assert error is None
        finally:
            file_path.unlink()
    
    @pytest.mark.asyncio
    async def test_validate_file_not_found(self):
        """Test validating a non-existent file."""
        processor = PowerPointProcessor()
        
        is_valid, error = await processor.validate_file(Path("/nonexistent/file.pptx"))
        
        assert is_valid is False
        assert "not found" in error.lower()
    
    @pytest.mark.asyncio
    async def test_validate_file_wrong_extension(self):
        """Test validating a file with wrong extension."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as f:
            f.write(b"test")
            file_path = Path(f.name)
        
        try:
            is_valid, error = await processor.validate_file(file_path)
            assert is_valid is False
            assert "unsupported" in error.lower()
        finally:
            file_path.unlink()
    
    @pytest.mark.asyncio
    async def test_validate_file_corrupted(self):
        """Test validating a corrupted file."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            f.write(b"not a valid pptx file")
            file_path = Path(f.name)
        
        try:
            is_valid, error = await processor.validate_file(file_path)
            assert is_valid is False
            assert "corrupted" in error.lower() or "cannot be read" in error.lower()
        finally:
            file_path.unlink()
    
    @pytest.mark.asyncio
    async def test_validate_file_empty_presentation(self):
        """Test validating a presentation with no slides."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            prs.save(f.name)
            file_path = Path(f.name)
        
        try:
            is_valid, error = await processor.validate_file(file_path)
            assert is_valid is False
            assert "no slides" in error.lower()
        finally:
            file_path.unlink()


class TestPowerPointProcessorWriteTranslated:
    """Tests for PowerPointProcessor write_translated method."""
    
    @pytest.mark.asyncio
    async def test_write_translated_shapes(self):
        """Test writing translated text to shapes."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Pt(100), Pt(100), Pt(200), Pt(100))
            txBox.text_frame.text = "Hello"
            prs.save(f.name)
            input_path = Path(f.name)
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)
        
        try:
            # Extract text
            segments = await processor.extract_text(input_path)
            shape_segments = [s for s in segments if s.metadata.get("type") == "shape"]
            
            # Create translations
            translations = ["Xin chào"]
            
            # Write translated (with output_mode="replace" to replace original text)
            success = await processor.write_translated(
                input_path, shape_segments, translations, output_path, output_mode="replace"
            )
            
            assert success is True
            assert output_path.exists()
            
            # Verify the translated content
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            assert "Xin chào" in texts
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()
    
    @pytest.mark.asyncio
    async def test_write_translated_preserves_bold(self):
        """Test that writing preserves bold formatting."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Pt(100), Pt(100), Pt(200), Pt(100))
            p = txBox.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Bold Text"
            run.font.bold = True
            prs.save(f.name)
            input_path = Path(f.name)
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(input_path)
            shape_segments = [s for s in segments if s.metadata.get("type") == "shape"]
            translations = ["Văn bản đậm"]
            
            # Write translated (with output_mode="replace" to replace original text)
            success = await processor.write_translated(
                input_path, shape_segments, translations, output_path, output_mode="replace"
            )
            
            assert success is True
            
            # Verify bold is preserved
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.runs and para.runs[0].text == "Văn bản đậm":
                            assert para.runs[0].font.bold is True
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()
    
    @pytest.mark.asyncio
    async def test_write_translated_table_cells(self):
        """Test writing translated text to table cells."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            table = slide.shapes.add_table(1, 2, Pt(100), Pt(100), Pt(400), Pt(100)).table
            table.cell(0, 0).text = "Hello"
            table.cell(0, 1).text = "World"
            prs.save(f.name)
            input_path = Path(f.name)
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(input_path)
            table_segments = [s for s in segments if s.metadata.get("type") == "table_cell"]
            
            translations = ["Xin chào", "Thế giới"]
            
            # Write translated (with output_mode="replace" to replace original text)
            success = await processor.write_translated(
                input_path, table_segments, translations, output_path, output_mode="replace"
            )
            
            assert success is True
            
            # Verify table content
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    assert table.cell(0, 0).text == "Xin chào"
                    assert table.cell(0, 1).text == "Thế giới"
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()
    
    @pytest.mark.asyncio
    async def test_write_translated_notes(self):
        """Test writing translated text to speaker notes."""
        processor = PowerPointProcessor()
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Original notes"
            prs.save(f.name)
            input_path = Path(f.name)
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)
        
        try:
            segments = await processor.extract_text(input_path)
            notes_segments = [s for s in segments if s.metadata.get("type") == "notes"]
            
            translations = ["Ghi chú đã dịch"]
            
            # Write translated (with output_mode="replace" to replace original text)
            success = await processor.write_translated(
                input_path, notes_segments, translations, output_path, output_mode="replace"
            )
            
            assert success is True
            
            # Verify notes content
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            notes_text = slide.notes_slide.notes_text_frame.text
            assert "Ghi chú đã dịch" in notes_text
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()


class TestPowerPointShapeIndexBug:
    """Tests that write-back targets the correct shape on mixed-content slides."""

    @staticmethod
    def _make_png():
        """Create a minimal 1x1 PNG image."""
        import struct, zlib
        sig = b'\x89PNG\r\n\x1a\n'
        ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
        ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
        raw = b'\x00\x00\x00\x00'
        idat_data = zlib.compress(raw)
        idat_crc = zlib.crc32(b'IDAT' + idat_data) & 0xffffffff
        idat = struct.pack('>I', len(idat_data)) + b'IDAT' + idat_data + struct.pack('>I', idat_crc)
        iend_crc = zlib.crc32(b'IEND') & 0xffffffff
        iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
        return sig + ihdr + idat + iend

    @pytest.mark.asyncio
    async def test_write_translated_mixed_shapes_targets_correct_shape(self):
        """On a slide with an image before a text box, translation should go to the text box."""
        processor = PowerPointProcessor()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # Shape 0: a picture (non-text shape)
            slide.shapes.add_picture(
                BytesIO(self._make_png()), Inches(0), Inches(0), Inches(1), Inches(1)
            )

            # Shape 1: a text box
            txBox = slide.shapes.add_textbox(Inches(2), Inches(0), Inches(3), Inches(1))
            txBox.text_frame.text = "Translate me"

            prs.save(f.name)
            input_path = Path(f.name)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)

        try:
            segments = await processor.extract_text(input_path)
            shape_segments = [s for s in segments if s.metadata.get("type") == "shape"]
            assert len(shape_segments) == 1
            assert shape_segments[0].text == "Translate me"
            # shape_idx should be 1 (second shape on slide), not 0
            assert shape_segments[0].metadata["shape_idx"] == 1

            translations = ["Dịch tôi"]
            success = await processor.write_translated(
                input_path, shape_segments, translations, output_path, output_mode="replace"
            )
            assert success is True

            # Verify text was written to the correct shape
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            shapes = list(slide.shapes)
            assert shapes[1].has_text_frame
            assert shapes[1].text_frame.text == "Dịch tôi"
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()

    @pytest.mark.asyncio
    async def test_write_translated_table_after_textbox(self):
        """Table write-back should target correct shape when preceded by a text box."""
        processor = PowerPointProcessor()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            # Shape 0: text box
            txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
            txBox.text_frame.text = "Header"

            # Shape 1: table
            table = slide.shapes.add_table(
                1, 1, Inches(0), Inches(2), Inches(3), Inches(1)
            ).table
            table.cell(0, 0).text = "Cell Data"

            prs.save(f.name)
            input_path = Path(f.name)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)

        try:
            segments = await processor.extract_text(input_path)
            table_segments = [s for s in segments if s.metadata.get("type") == "table_cell"]
            assert len(table_segments) == 1
            assert table_segments[0].metadata["shape_idx"] == 1

            translations = ["Dữ liệu ô"]
            success = await processor.write_translated(
                input_path, table_segments, translations, output_path, output_mode="replace"
            )
            assert success is True

            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            shapes = list(slide.shapes)
            assert shapes[1].has_table
            assert shapes[1].table.cell(0, 0).text == "Dữ liệu ô"
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()


class TestPowerPointChartExtraction:
    """Tests for chart text extraction and write-back."""

    @pytest.mark.asyncio
    async def test_extract_chart_title(self):
        """Test extracting text from a chart title."""
        processor = PowerPointProcessor()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            chart_data = CategoryChartData()
            chart_data.categories = ['East', 'West']
            chart_data.add_series('Sales', (10, 20))

            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1), Inches(6), Inches(4),
                chart_data
            )
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.has_text_frame = True
            chart_shape.chart.chart_title.text_frame.text = "Sales Report"

            prs.save(f.name)
            input_path = Path(f.name)

        try:
            segments = await processor.extract_text(input_path)
            chart_segments = [s for s in segments if s.metadata.get("type") == "chart"]
            assert len(chart_segments) >= 1

            title_segments = [
                s for s in chart_segments if s.metadata.get("chart_element") == "title"
            ]
            assert len(title_segments) >= 1
            assert "Sales Report" in title_segments[0].text
        finally:
            input_path.unlink()

    @pytest.mark.asyncio
    async def test_write_translated_chart_title(self):
        """Test writing translated chart title."""
        processor = PowerPointProcessor()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            chart_data = CategoryChartData()
            chart_data.categories = ['East', 'West']
            chart_data.add_series('Sales', (10, 20))

            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1), Inches(6), Inches(4),
                chart_data
            )
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.has_text_frame = True
            chart_shape.chart.chart_title.text_frame.text = "Sales Report"

            prs.save(f.name)
            input_path = Path(f.name)

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            output_path = Path(f.name)

        try:
            segments = await processor.extract_text(input_path)
            chart_segments = [s for s in segments if s.metadata.get("type") == "chart"]
            translations = ["Báo cáo bán hàng"] * len(chart_segments)

            success = await processor.write_translated(
                input_path, chart_segments, translations, output_path, output_mode="replace"
            )
            assert success is True

            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.has_chart and shape.chart.has_title:
                    assert "Báo cáo bán hàng" in shape.chart.chart_title.text_frame.text
        finally:
            input_path.unlink()
            if output_path.exists():
                output_path.unlink()

    @pytest.mark.asyncio
    async def test_extract_chart_with_axis_titles(self):
        """Test extracting axis titles from a chart."""
        processor = PowerPointProcessor()

        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)

            chart_data = CategoryChartData()
            chart_data.categories = ['Q1', 'Q2']
            chart_data.add_series('Revenue', (100, 200))

            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(1), Inches(6), Inches(4),
                chart_data
            )
            chart = chart_shape.chart

            # Add value axis title
            value_axis = chart.value_axis
            value_axis.has_title = True
            value_axis.axis_title.has_text_frame = True
            value_axis.axis_title.text_frame.text = "Revenue Amount"

            # Add category axis title
            category_axis = chart.category_axis
            category_axis.has_title = True
            category_axis.axis_title.has_text_frame = True
            category_axis.axis_title.text_frame.text = "Quarter"

            prs.save(f.name)
            input_path = Path(f.name)

        try:
            segments = await processor.extract_text(input_path)
            chart_segments = [s for s in segments if s.metadata.get("type") == "chart"]

            elements = {s.metadata["chart_element"]: s.text for s in chart_segments}
            assert "value_axis_title" in elements
            assert "Revenue Amount" in elements["value_axis_title"]
            assert "category_axis_title" in elements
            assert "Quarter" in elements["category_axis_title"]
        finally:
            input_path.unlink()
