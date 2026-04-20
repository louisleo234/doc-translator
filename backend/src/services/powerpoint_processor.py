"""
PowerPoint Document Processor Module

Provides document processing for Microsoft PowerPoint (.pptx) files.
Extracts text from slides, notes, text boxes, shapes, and tables
while preserving formatting during translation writing.
"""

import logging
from pathlib import Path
from typing import List, Optional
import asyncio

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table

from .document_processor import (
    DocumentProcessor,
    DocumentType,
    TextSegment,
    apply_output_mode,
)


class PowerPointProcessor(DocumentProcessor):
    """
    Document processor for Microsoft PowerPoint files (.pptx).
    
    Extracts text from:
    - Slide text frames and shapes
    - Tables on slides
    - Speaker notes
    - Master slides and layouts
    
    Preserves formatting during translation writing:
    - Font properties (name, size, color, bold, italic, underline)
    - Shape properties and positions
    - Slide layouts and themes
    - Animations and transitions
    - Embedded media
    """
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """
        Initialize the PowerPoint document processor.
        
        Args:
            logger: Logger instance for logging operations
        """
        self.logger = logger or logging.getLogger(__name__)

    @property
    def supported_extensions(self) -> List[str]:
        """Return list of supported file extensions."""
        return ['.pptx']
    
    @property
    def document_type(self) -> DocumentType:
        """Return the document type this processor handles."""
        return DocumentType.POWERPOINT

    async def extract_text(self, file_path: Path) -> List[TextSegment]:
        """
        Extract all translatable text segments from the PowerPoint presentation.
        
        Extracts text from:
        - Slide shapes (text boxes, placeholders)
        - Tables on slides
        - Speaker notes
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of TextSegment objects containing text and metadata
            
        Raises:
            ValueError: If the file cannot be loaded
        """
        try:
            presentation = await asyncio.to_thread(Presentation, str(file_path))
        except Exception as e:
            error_msg = str(e).lower()
            if "package" in error_msg or "zip" in error_msg:
                raise ValueError(f"Failed to load PowerPoint file: {file_path}. File may be corrupted.")
            raise ValueError(f"Failed to load PowerPoint file: {file_path}. Error: {str(e)}")
        
        segments: List[TextSegment] = []
        segment_id = 0
        
        # Extract from each slide
        for slide_idx, slide in enumerate(presentation.slides):
            # Extract from shapes on the slide
            shape_segments = self._extract_shape_segments(
                slide.shapes, slide_idx, segment_id
            )
            segments.extend(shape_segments)
            segment_id += len(shape_segments)
            
            # Extract from speaker notes
            notes_segments = self._extract_notes_segments(
                slide, slide_idx, segment_id
            )
            segments.extend(notes_segments)
            segment_id += len(notes_segments)
        
        self.logger.info(f"Extracted {len(segments)} text segments from {file_path.name}")
        return segments
    
    def _extract_shape_segments(
        self,
        shapes,
        slide_idx: int,
        start_segment_id: int
    ) -> List[TextSegment]:
        """Extract text segments from shapes on a slide."""
        segments = []
        segment_id = start_segment_id
        
        for shape_idx, shape in enumerate(shapes):
            # Handle grouped shapes recursively
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_segments = self._extract_group_shape_segments(
                    shape, slide_idx, shape_idx, segment_id
                )
                segments.extend(group_segments)
                segment_id += len(group_segments)
                continue
            
            # Handle tables
            if shape.has_table:
                table_segments = self._extract_table_segments(
                    shape.table, slide_idx, shape_idx, segment_id
                )
                segments.extend(table_segments)
                segment_id += len(table_segments)
                continue

            # Handle charts
            if shape.has_chart:
                chart_segments = self._extract_chart_segments(
                    shape.chart, slide_idx, shape_idx, segment_id
                )
                segments.extend(chart_segments)
                segment_id += len(chart_segments)
                continue

            # Handle shapes with text frames
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    para_text = paragraph.text.strip()
                    if para_text:
                        segment = self._create_shape_segment(
                            shape, paragraph, slide_idx, shape_idx, para_idx, segment_id
                        )
                        segments.append(segment)
                        segment_id += 1
        
        return segments
    
    def _extract_group_shape_segments(
        self,
        group_shape,
        slide_idx: int,
        group_idx: int,
        start_segment_id: int
    ) -> List[TextSegment]:
        """Extract text segments from a group of shapes, including nested groups, tables, and charts."""
        segments = []
        segment_id = start_segment_id

        for shape_idx, shape in enumerate(group_shape.shapes):
            # Handle nested groups recursively
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                nested_segments = self._extract_group_shape_segments(
                    shape, slide_idx, group_idx, segment_id
                )
                # Tag nested segments with the outer group context
                for seg in nested_segments:
                    seg.metadata["nested_group_idx"] = shape_idx
                segments.extend(nested_segments)
                segment_id += len(nested_segments)
                continue

            # Handle tables within groups
            if shape.has_table:
                table_segments = self._extract_table_segments(
                    shape.table, slide_idx, shape_idx, segment_id
                )
                # Tag with group context
                for seg in table_segments:
                    seg.metadata["group_idx"] = group_idx
                    seg.metadata["type"] = "group_table_cell"
                segments.extend(table_segments)
                segment_id += len(table_segments)
                continue

            # Handle charts within groups
            if shape.has_chart:
                chart_segments = self._extract_chart_segments(
                    shape.chart, slide_idx, shape_idx, segment_id
                )
                for seg in chart_segments:
                    seg.metadata["group_idx"] = group_idx
                segments.extend(chart_segments)
                segment_id += len(chart_segments)
                continue

            if shape.has_text_frame:
                text_frame = shape.text_frame
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    para_text = paragraph.text.strip()
                    if para_text:
                        runs_metadata = self._get_runs_metadata(paragraph)

                        segment = TextSegment(
                            id=str(segment_id),
                            text=para_text,
                            location=f"Slide {slide_idx + 1}, Group {group_idx + 1}, Shape {shape_idx + 1}, Paragraph {para_idx + 1}",
                            metadata={
                                "type": "group_shape",
                                "slide_idx": slide_idx,
                                "group_idx": group_idx,
                                "shape_idx": shape_idx,
                                "paragraph_idx": para_idx,
                                "runs": runs_metadata,
                            }
                        )
                        segments.append(segment)
                        segment_id += 1

        return segments
    
    def _create_shape_segment(
        self,
        shape,
        paragraph,
        slide_idx: int,
        shape_idx: int,
        para_idx: int,
        segment_id: int
    ) -> TextSegment:
        """Create a TextSegment from a shape paragraph with run-level metadata."""
        runs_metadata = self._get_runs_metadata(paragraph)
        
        # Determine shape type for location description
        shape_type_name = self._get_shape_type_name(shape)
        
        return TextSegment(
            id=str(segment_id),
            text=paragraph.text,
            location=f"Slide {slide_idx + 1}, {shape_type_name} {shape_idx + 1}, Paragraph {para_idx + 1}",
            metadata={
                "type": "shape",
                "slide_idx": slide_idx,
                "shape_idx": shape_idx,
                "paragraph_idx": para_idx,
                "shape_type": str(shape.shape_type) if hasattr(shape, 'shape_type') else None,
                "runs": runs_metadata,
            }
        )
    
    def _get_shape_type_name(self, shape) -> str:
        """Get a human-readable name for the shape type."""
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            return "Placeholder"
        if hasattr(shape, 'shape_type'):
            shape_type = shape.shape_type
            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return "Text Box"
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "Shape"
        return "Shape"
    
    def _get_runs_metadata(self, paragraph) -> List[dict]:
        """Extract metadata from paragraph runs for format preservation."""
        runs_metadata = []
        for run_idx, run in enumerate(paragraph.runs):
            run_meta = {
                "run_idx": run_idx,
                "text": run.text,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "underline": run.font.underline,
                "font_name": run.font.name,
                "font_size": run.font.size.pt if run.font.size else None,
                "font_color": self._get_font_color(run),
            }
            runs_metadata.append(run_meta)
        return runs_metadata
    
    def _get_font_color(self, run) -> Optional[str]:
        """Extract font color as hex string if available."""
        try:
            if run.font.color and run.font.color.rgb:
                return str(run.font.color.rgb)
        except Exception:
            pass
        return None
    
    def _extract_table_segments(
        self,
        table: Table,
        slide_idx: int,
        shape_idx: int,
        start_segment_id: int
    ) -> List[TextSegment]:
        """Extract text segments from a table on a slide."""
        segments = []
        segment_id = start_segment_id
        
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                if cell_text:
                    # Get paragraph-level metadata for the cell
                    cell_paragraphs_meta = []
                    for para_idx, para in enumerate(cell.text_frame.paragraphs):
                        if para.text.strip():
                            runs_meta = self._get_runs_metadata(para)
                            cell_paragraphs_meta.append({
                                "para_idx": para_idx,
                                "text": para.text,
                                "runs": runs_meta,
                            })
                    
                    segment = TextSegment(
                        id=str(segment_id),
                        text=cell_text,
                        location=f"Slide {slide_idx + 1}, Table {shape_idx + 1}, Row {row_idx + 1}, Column {col_idx + 1}",
                        metadata={
                            "type": "table_cell",
                            "slide_idx": slide_idx,
                            "shape_idx": shape_idx,
                            "row_idx": row_idx,
                            "col_idx": col_idx,
                            "paragraphs": cell_paragraphs_meta,
                        }
                    )
                    segments.append(segment)
                    segment_id += 1
        
        return segments
    
    def _extract_chart_segments(
        self,
        chart,
        slide_idx: int,
        shape_idx: int,
        start_segment_id: int
    ) -> List[TextSegment]:
        """Extract translatable text segments from a chart (title, axis titles)."""
        segments = []
        segment_id = start_segment_id

        # Chart title
        if chart.has_title and chart.chart_title.has_text_frame:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text:
                for para_idx, para in enumerate(chart.chart_title.text_frame.paragraphs):
                    para_text = para.text.strip()
                    if para_text:
                        runs_metadata = self._get_runs_metadata(para)
                        segment = TextSegment(
                            id=str(segment_id),
                            text=para_text,
                            location=f"Slide {slide_idx + 1}, Chart {shape_idx + 1}, Title, Paragraph {para_idx + 1}",
                            metadata={
                                "type": "chart",
                                "chart_element": "title",
                                "slide_idx": slide_idx,
                                "shape_idx": shape_idx,
                                "paragraph_idx": para_idx,
                                "runs": runs_metadata,
                            }
                        )
                        segments.append(segment)
                        segment_id += 1

        # Category axis title
        try:
            if chart.category_axis.has_title:
                axis_title = chart.category_axis.axis_title
                if axis_title.has_text_frame:
                    for para_idx, para in enumerate(axis_title.text_frame.paragraphs):
                        para_text = para.text.strip()
                        if para_text:
                            runs_metadata = self._get_runs_metadata(para)
                            segment = TextSegment(
                                id=str(segment_id),
                                text=para_text,
                                location=f"Slide {slide_idx + 1}, Chart {shape_idx + 1}, Category Axis Title, Paragraph {para_idx + 1}",
                                metadata={
                                    "type": "chart",
                                    "chart_element": "category_axis_title",
                                    "slide_idx": slide_idx,
                                    "shape_idx": shape_idx,
                                    "paragraph_idx": para_idx,
                                    "runs": runs_metadata,
                                }
                            )
                            segments.append(segment)
                            segment_id += 1
        except (ValueError, AttributeError):
            pass

        # Value axis title
        try:
            if chart.value_axis.has_title:
                axis_title = chart.value_axis.axis_title
                if axis_title.has_text_frame:
                    for para_idx, para in enumerate(axis_title.text_frame.paragraphs):
                        para_text = para.text.strip()
                        if para_text:
                            runs_metadata = self._get_runs_metadata(para)
                            segment = TextSegment(
                                id=str(segment_id),
                                text=para_text,
                                location=f"Slide {slide_idx + 1}, Chart {shape_idx + 1}, Value Axis Title, Paragraph {para_idx + 1}",
                                metadata={
                                    "type": "chart",
                                    "chart_element": "value_axis_title",
                                    "slide_idx": slide_idx,
                                    "shape_idx": shape_idx,
                                    "paragraph_idx": para_idx,
                                    "runs": runs_metadata,
                                }
                            )
                            segments.append(segment)
                            segment_id += 1
        except (ValueError, AttributeError):
            pass

        return segments

    def _extract_notes_segments(
        self,
        slide,
        slide_idx: int,
        start_segment_id: int
    ) -> List[TextSegment]:
        """Extract text segments from speaker notes."""
        segments = []
        segment_id = start_segment_id
        
        if not slide.has_notes_slide:
            return segments
        
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        if notes_text_frame:
            for para_idx, paragraph in enumerate(notes_text_frame.paragraphs):
                para_text = paragraph.text.strip()
                if para_text:
                    runs_metadata = self._get_runs_metadata(paragraph)
                    
                    segment = TextSegment(
                        id=str(segment_id),
                        text=para_text,
                        location=f"Slide {slide_idx + 1}, Notes, Paragraph {para_idx + 1}",
                        metadata={
                            "type": "notes",
                            "slide_idx": slide_idx,
                            "paragraph_idx": para_idx,
                            "runs": runs_metadata,
                        }
                    )
                    segments.append(segment)
                    segment_id += 1
        
        return segments


    async def write_translated(
        self,
        file_path: Path,
        segments: List[TextSegment],
        translations: List[str],
        output_path: Path,
        output_mode: str = "replace"
    ) -> bool:
        """
        Write translated text back to PowerPoint presentation, preserving formatting.

        Args:
            file_path: Path to the original PowerPoint file
            segments: List of original text segments
            translations: List of translated texts (same order as segments)
            output_path: Path where the translated presentation should be saved
            output_mode: One of "replace", "append", "prepend", "interleave", "interleave_reverse" (default: "replace")
            
        Returns:
            True if writing succeeded, False otherwise
        """
        try:
            # Always load fresh from file to avoid shared state issues
            # during concurrent processing of multiple files
            presentation = await asyncio.to_thread(Presentation, str(file_path))
            
            # Create translation map with output mode applied
            translation_map = {}
            for seg, trans in zip(segments, translations):
                final_text = apply_output_mode(seg.text, trans, output_mode)
                translation_map[seg.id] = final_text
            
            # Process each segment
            for segment in segments:
                translation = translation_map.get(segment.id)
                if translation is None:
                    continue
                
                seg_type = segment.metadata.get("type")
                
                if seg_type == "shape":
                    await self._write_shape_translation(
                        presentation, segment, translation
                    )
                elif seg_type == "group_shape":
                    await self._write_group_shape_translation(
                        presentation, segment, translation
                    )
                elif seg_type == "table_cell":
                    await self._write_table_cell_translation(
                        presentation, segment, translation
                    )
                elif seg_type == "chart":
                    await self._write_chart_translation(
                        presentation, segment, translation
                    )
                elif seg_type == "group_table_cell":
                    await self._write_table_cell_translation(
                        presentation, segment, translation
                    )
                elif seg_type == "notes":
                    await self._write_notes_translation(
                        presentation, segment, translation
                    )
            
            # Ensure output directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save the presentation
            await asyncio.to_thread(presentation.save, str(output_path))
            
            self.logger.info(f"Saved translated PowerPoint presentation to: {output_path}")
            return True
            
        except Exception as e:
            self.logger.error(f"Failed to write translated PowerPoint presentation: {e}")
            return False

    async def _write_shape_translation(
        self,
        presentation: Presentation,
        segment: TextSegment,
        translation: str
    ) -> None:
        """Write translation to a shape while preserving formatting."""
        slide_idx = segment.metadata.get("slide_idx", 0)
        shape_idx = segment.metadata.get("shape_idx", 0)
        para_idx = segment.metadata.get("paragraph_idx", 0)
        
        if slide_idx >= len(presentation.slides):
            self.logger.warning(f"Slide index {slide_idx} out of range")
            return
        
        slide = presentation.slides[slide_idx]
        shapes = list(slide.shapes)

        if shape_idx >= len(shapes):
            self.logger.warning(f"Shape index {shape_idx} out of range on slide {slide_idx}")
            return

        shape = shapes[shape_idx]
        if not shape.has_text_frame:
            self.logger.warning(f"Shape {shape_idx} on slide {slide_idx} has no text frame")
            return

        text_frame = shape.text_frame
        
        if para_idx >= len(text_frame.paragraphs):
            self.logger.warning(f"Paragraph index {para_idx} out of range in shape {shape_idx}")
            return
        
        paragraph = text_frame.paragraphs[para_idx]
        runs_meta = segment.metadata.get("runs", [])
        
        self._update_paragraph_text(paragraph, translation, runs_meta)

    async def _write_group_shape_translation(
        self,
        presentation: Presentation,
        segment: TextSegment,
        translation: str
    ) -> None:
        """Write translation to a shape within a group."""
        slide_idx = segment.metadata.get("slide_idx", 0)
        group_idx = segment.metadata.get("group_idx", 0)
        shape_idx = segment.metadata.get("shape_idx", 0)
        para_idx = segment.metadata.get("paragraph_idx", 0)
        
        if slide_idx >= len(presentation.slides):
            self.logger.warning(f"Slide index {slide_idx} out of range")
            return
        
        slide = presentation.slides[slide_idx]
        shapes = list(slide.shapes)

        if group_idx >= len(shapes):
            self.logger.warning(f"Group index {group_idx} out of range on slide {slide_idx}")
            return

        group_shape = shapes[group_idx]
        if group_shape.shape_type != MSO_SHAPE_TYPE.GROUP:
            self.logger.warning(f"Shape {group_idx} on slide {slide_idx} is not a group")
            return

        inner_shapes = list(group_shape.shapes)
        if shape_idx >= len(inner_shapes):
            self.logger.warning(f"Shape index {shape_idx} out of range in group {group_idx}")
            return

        shape = inner_shapes[shape_idx]
        if not shape.has_text_frame:
            self.logger.warning(f"Shape {shape_idx} in group {group_idx} has no text frame")
            return
        text_frame = shape.text_frame
        
        if para_idx >= len(text_frame.paragraphs):
            self.logger.warning(f"Paragraph index {para_idx} out of range")
            return
        
        paragraph = text_frame.paragraphs[para_idx]
        runs_meta = segment.metadata.get("runs", [])
        
        self._update_paragraph_text(paragraph, translation, runs_meta)

    async def _write_table_cell_translation(
        self,
        presentation: Presentation,
        segment: TextSegment,
        translation: str
    ) -> None:
        """Write translation to a table cell while preserving formatting."""
        slide_idx = segment.metadata.get("slide_idx", 0)
        shape_idx = segment.metadata.get("shape_idx", 0)
        row_idx = segment.metadata.get("row_idx", 0)
        col_idx = segment.metadata.get("col_idx", 0)
        
        if slide_idx >= len(presentation.slides):
            self.logger.warning(f"Slide index {slide_idx} out of range")
            return
        
        slide = presentation.slides[slide_idx]
        shapes = list(slide.shapes)

        if shape_idx >= len(shapes):
            self.logger.warning(f"Table shape index {shape_idx} out of range on slide {slide_idx}")
            return

        shape = shapes[shape_idx]
        if not shape.has_table:
            self.logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a table")
            return

        table = shape.table
        
        if row_idx >= len(table.rows):
            self.logger.warning(f"Row index {row_idx} out of range in table")
            return
        
        row = table.rows[row_idx]
        
        if col_idx >= len(row.cells):
            self.logger.warning(f"Column index {col_idx} out of range in table")
            return
        
        cell = row.cells[col_idx]
        
        # Get paragraph metadata if available
        paragraphs_meta = segment.metadata.get("paragraphs", [])
        
        if cell.text_frame.paragraphs:
            # Update the first paragraph with the translation
            first_para = cell.text_frame.paragraphs[0]
            
            if paragraphs_meta and paragraphs_meta[0].get("runs"):
                runs_meta = paragraphs_meta[0]["runs"]
                self._update_paragraph_text(first_para, translation, runs_meta)
            else:
                self._update_paragraph_text_simple(first_para, translation)
            
            # Clear any additional paragraphs
            for para in cell.text_frame.paragraphs[1:]:
                self._clear_paragraph(para)

    async def _write_notes_translation(
        self,
        presentation: Presentation,
        segment: TextSegment,
        translation: str
    ) -> None:
        """Write translation to speaker notes."""
        slide_idx = segment.metadata.get("slide_idx", 0)
        para_idx = segment.metadata.get("paragraph_idx", 0)
        
        if slide_idx >= len(presentation.slides):
            self.logger.warning(f"Slide index {slide_idx} out of range")
            return
        
        slide = presentation.slides[slide_idx]
        
        if not slide.has_notes_slide:
            self.logger.warning(f"Slide {slide_idx} has no notes slide")
            return
        
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        if para_idx >= len(notes_text_frame.paragraphs):
            self.logger.warning(f"Paragraph index {para_idx} out of range in notes")
            return
        
        paragraph = notes_text_frame.paragraphs[para_idx]
        runs_meta = segment.metadata.get("runs", [])
        
        self._update_paragraph_text(paragraph, translation, runs_meta)

    async def _write_chart_translation(
        self,
        presentation: Presentation,
        segment: TextSegment,
        translation: str
    ) -> None:
        """Write translation to a chart element (title or axis title)."""
        slide_idx = segment.metadata.get("slide_idx", 0)
        shape_idx = segment.metadata.get("shape_idx", 0)
        chart_element = segment.metadata.get("chart_element", "")
        para_idx = segment.metadata.get("paragraph_idx", 0)

        if slide_idx >= len(presentation.slides):
            self.logger.warning(f"Slide index {slide_idx} out of range")
            return

        slide = presentation.slides[slide_idx]
        shapes = list(slide.shapes)

        if shape_idx >= len(shapes):
            self.logger.warning(f"Shape index {shape_idx} out of range on slide {slide_idx}")
            return

        shape = shapes[shape_idx]
        if not shape.has_chart:
            self.logger.warning(f"Shape {shape_idx} on slide {slide_idx} is not a chart")
            return

        chart = shape.chart
        text_frame = None

        if chart_element == "title" and chart.has_title and chart.chart_title.has_text_frame:
            text_frame = chart.chart_title.text_frame
        elif chart_element == "category_axis_title":
            try:
                if chart.category_axis.has_title:
                    axis_title = chart.category_axis.axis_title
                    if axis_title.has_text_frame:
                        text_frame = axis_title.text_frame
            except (ValueError, AttributeError):
                pass
        elif chart_element == "value_axis_title":
            try:
                if chart.value_axis.has_title:
                    axis_title = chart.value_axis.axis_title
                    if axis_title.has_text_frame:
                        text_frame = axis_title.text_frame
            except (ValueError, AttributeError):
                pass

        if text_frame is None:
            self.logger.warning(f"Chart element '{chart_element}' not found on slide {slide_idx}")
            return

        if para_idx >= len(text_frame.paragraphs):
            self.logger.warning(f"Paragraph index {para_idx} out of range in chart element")
            return

        paragraph = text_frame.paragraphs[para_idx]
        runs_meta = segment.metadata.get("runs", [])
        self._update_paragraph_text(paragraph, translation, runs_meta)

    def _update_paragraph_text(
        self,
        paragraph,
        translation: str,
        runs_meta: List[dict]
    ) -> None:
        """Update paragraph text while preserving run-level formatting."""
        if not runs_meta:
            self._update_paragraph_text_simple(paragraph, translation)
            return
        
        # Get formatting from first run
        first_run_meta = runs_meta[0]
        
        # Clear text from all runs except the first (can't remove runs, only clear them)
        runs = list(paragraph.runs)
        for i, run in enumerate(runs):
            if i > 0:
                run.text = ""
        
        if paragraph.runs:
            # Update the first run with translation
            run = paragraph.runs[0]
            run.text = translation
            
            # Apply preserved formatting
            if first_run_meta.get("bold") is not None:
                run.font.bold = first_run_meta["bold"]
            if first_run_meta.get("italic") is not None:
                run.font.italic = first_run_meta["italic"]
            if first_run_meta.get("underline") is not None:
                run.font.underline = first_run_meta["underline"]
            if first_run_meta.get("font_name"):
                run.font.name = first_run_meta["font_name"]
            if first_run_meta.get("font_size"):
                run.font.size = Pt(first_run_meta["font_size"])
            if first_run_meta.get("font_color"):
                try:
                    color_str = first_run_meta["font_color"]
                    run.font.color.rgb = RGBColor.from_string(color_str)
                except Exception:
                    pass
        else:
            # No runs exist, add a new one
            run = paragraph.add_run()
            run.text = translation

    def _update_paragraph_text_simple(self, paragraph, translation: str) -> None:
        """Update paragraph text when there's a single run or simple structure."""
        if paragraph.runs:
            # Store formatting from first run
            first_run = paragraph.runs[0]
            bold = first_run.font.bold
            italic = first_run.font.italic
            underline = first_run.font.underline
            font_name = first_run.font.name
            font_size = first_run.font.size
            
            # Clear text from all runs except the first (can't remove runs, only clear them)
            runs = list(paragraph.runs)
            for i, run in enumerate(runs):
                if i > 0:
                    run.text = ""
            
            # Update first run
            first_run.text = translation
            first_run.font.bold = bold
            first_run.font.italic = italic
            first_run.font.underline = underline
            if font_name:
                first_run.font.name = font_name
            if font_size:
                first_run.font.size = font_size
        else:
            run = paragraph.add_run()
            run.text = translation

    def _clear_paragraph(self, paragraph) -> None:
        """Clear all text from a paragraph."""
        for run in paragraph.runs:
            run.text = ""


    async def validate_file(self, file_path: Path) -> tuple[bool, Optional[str]]:
        """
        Validate that the PowerPoint file can be processed.
        
        Checks for:
        - File existence
        - File readability (not password protected or corrupted)
        - Valid .pptx format
        
        Args:
            file_path: Path to the PowerPoint file to validate
            
        Returns:
            Tuple of (is_valid, error_message)
        """
        if not file_path.exists():
            return False, f"File not found: {file_path}"
        
        if file_path.suffix.lower() not in self.supported_extensions:
            return False, f"Unsupported file format: {file_path.suffix}"
        
        try:
            # Try to load the presentation to check for corruption/password protection
            presentation = await asyncio.to_thread(Presentation, str(file_path))
            
            # Check if presentation has any slides
            if len(presentation.slides) == 0:
                return False, "PowerPoint presentation contains no slides."
            
            # Check if there's any text content
            has_content = False
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                has_content = True
                                break
                    if shape.has_table:
                        has_content = True
                        break
                    if shape.has_chart:
                        has_content = True
                        break
                    if has_content:
                        break
                if has_content:
                    break
            
            if not has_content:
                return False, "PowerPoint presentation contains no translatable text content."
            
            return True, None
            
        except Exception as e:
            error_msg = str(e).lower()
            if "password" in error_msg or "encrypted" in error_msg:
                return False, "This file is password protected. Please remove the password and try again."
            if "package" in error_msg or "zip" in error_msg or "bad" in error_msg:
                return False, "This file appears to be corrupted and cannot be read."
            return False, f"Failed to read PowerPoint file: {str(e)}"
